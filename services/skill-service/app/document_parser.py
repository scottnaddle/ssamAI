"""Document parsing: extract structured content from .hwpx, .pptx files.

Supported formats:
- .hwpx (Hancom Word XML): ZIP archive → Contents/content.hpl XML parsing
- .pptx (PowerPoint): python-pptx → slide-by-slide text extraction
- .hwp (HWP 5.0 binary): basic text extraction via olefile (experimental)
"""

from __future__ import annotations

import re
import zipfile
from dataclasses import dataclass, field
from io import BytesIO
from pathlib import Path
from typing import Any
from xml.etree import ElementTree as ET

# ── OLE (old .hwp format) — optional dependency ──────────────
try:
    import olefile
    HAS_OLEFILE = True
except ImportError:
    HAS_OLEFILE = False

# ── python-pptx (.pptx) — optional dependency ─────────────────
try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False


@dataclass
class ParsedSection:
    """A logical section of a parsed document (e.g. learning objectives, teaching flow)."""
    heading: str | None = None
    body_lines: list[str] = field(default_factory=list)
    table_rows: list[list[str]] = field(default_factory=list)


@dataclass
class ParsedDocument:
    """Structured output of document parsing."""
    title: str | None = None
    sections: list[ParsedSection] = field(default_factory=list)
    full_text: str = ""
    raw_metadata: dict[str, Any] = field(default_factory=dict)


# ─── HWPX Parser ─────────────────────────────────────────────

_HWPX_NS = "http://www.hancom.co.kr/hml/2021"
_HWPX_NS_MAP = {"h": _HWPX_NS}


def _hwpx_extract_text(element: ET.Element) -> str:
    """Recursively extract text from an HWPX XML element."""
    texts: list[str] = []
    # Text spans inside <h:t> elements
    for t in element.iterfind(".//h:t", _HWPX_NS_MAP):
        text = (t.text or "").strip()
        if text:
            texts.append(text)
    # Fallback: direct text content
    if not texts and element.text:
        texts.append(element.text.strip())
    return " ".join(texts)


def _parse_hwpx_sections(root: ET.Element) -> list[ParsedSection]:
    """Extract sections from HWPX content XML."""
    sections: list[ParsedSection] = []

    for section_elem in root.iterfind(".//h:section", _HWPX_NS_MAP):
        sec = ParsedSection()
        # Try to find section heading
        for p in section_elem.iterfind(".//h:p", _HWPX_NS_MAP):
            text = _hwpx_extract_text(p)
            if text:
                # Heuristic: short first paragraph is likely a heading
                if sec.heading is None and len(text) < 80:
                    sec.heading = text
                else:
                    sec.body_lines.append(text)
        if sec.heading or sec.body_lines:
            sections.append(sec)

    return sections


def parse_hwpx(data: bytes) -> ParsedDocument:
    """Parse a .hwpx file (ZIP of XML) into structured document."""
    doc = ParsedDocument()
    try:
        with zipfile.ZipFile(BytesIO(data)) as zf:
            content_file = None
            for name in zf.namelist():
                if name.endswith("content.hpl") or name.endswith("content.xml"):
                    content_file = name
                    break
            if content_file is None:
                doc.full_text = "HWPX 파일에서 content.hpl을 찾을 수 없습니다."
                return doc

            xml_bytes = zf.read(content_file)
            root = ET.fromstring(xml_bytes)

            # Extract full text
            all_texts: list[str] = []
            for elem in root.iter():
                if elem.text and elem.text.strip():
                    all_texts.append(elem.text.strip())
            doc.full_text = "\n".join(all_texts)

            # Extract sections
            doc.sections = _parse_hwpx_sections(root)

            # Try to find title from document properties or first heading
            for p in root.iterfind(".//h:p", _HWPX_NS_MAP):
                text = _hwpx_extract_text(p)
                if text:
                    doc.title = text[:100]
                    break

    except Exception as exc:
        doc.full_text = f"HWPX 파싱 오류: {exc}"

    return doc


# ─── PPTX Parser ─────────────────────────────────────────────

def parse_pptx(data: bytes) -> ParsedDocument:
    """Parse a .pptx file into structured document (one section per slide)."""
    doc = ParsedDocument()
    if not HAS_PPTX:
        doc.full_text = "python-pptx가 설치되지 않아 .pptx 파싱이 불가능합니다."
        return doc

    try:
        prs = Presentation(BytesIO(data))
        all_texts: list[str] = []
        for i, slide in enumerate(prs.slides, 1):
            sec = ParsedSection(heading=f"슬라이드 {i}")
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            sec.body_lines.append(text)
                            all_texts.append(text)
                if shape.has_table:
                    table = shape.table
                    rows: list[list[str]] = []
                    for row in table.rows:
                        row_texts = [cell.text.strip() for cell in row.cells]
                        if any(row_texts):
                            rows.append(row_texts)
                    if rows:
                        sec.table_rows = rows
            if sec.body_lines or sec.table_rows:
                doc.sections.append(sec)

        doc.full_text = "\n".join(all_texts)
        doc.title = doc.sections[0].heading if doc.sections else None
    except Exception as exc:
        doc.full_text = f"PPTX 파싱 오류: {exc}"

    return doc


# ─── HWP Binary Parser (experimental) ─────────────────────────

def parse_hwp(data: bytes) -> ParsedDocument:
    """Basic text extraction from HWP 5.0 binary format using olefile."""
    doc = ParsedDocument()
    if not HAS_OLEFILE:
        doc.full_text = (
            "HWP 5.0 파싱을 위해 olefile 패키지가 필요합니다. "
            "pip install olefile 로 설치하거나 .hwpx 형식으로 변환해 업로드해주세요."
        )
        return doc

    try:
        ole = olefile.OleFileIO(BytesIO(data))
        # HWP stores main text in PrvText stream
        if ole.exists("PrvText"):
            text_bytes = ole.openstream("PrvText").read()
            try:
                text = text_bytes.decode("utf-16-le", errors="replace")
            except UnicodeDecodeError:
                text = text_bytes.decode("euc-kr", errors="replace")
            doc.full_text = text
            # Split into lines for basic structure
            lines = [l.strip() for l in text.split("\n") if l.strip()]
            if lines:
                doc.title = lines[0][:100]
                # Group remaining lines into one section
                doc.sections = [ParsedSection(body_lines=lines[1:])]
        else:
            doc.full_text = "HWP 파일에서 PrvText 스트림을 찾을 수 없습니다."

        # Also try PrvText.utf8 for newer HWP formats
        if not doc.full_text and ole.exists("PrvText.utf8"):
            text_bytes = ole.openstream("PrvText.utf8").read()
            doc.full_text = text_bytes.decode("utf-8", errors="replace")

        ole.close()
    except Exception as exc:
        doc.full_text = f"HWP 파싱 오류: {exc}"

    return doc


# ─── Dispatcher ──────────────────────────────────────────────

def parse_document(data: bytes, filename: str) -> ParsedDocument:
    """Detect format and parse accordingly."""
    ext = Path(filename).suffix.lower()
    if ext == ".hwpx":
        return parse_hwpx(data)
    elif ext == ".pptx":
        return parse_pptx(data)
    elif ext == ".hwp":
        return parse_hwp(data)
    else:
        doc = ParsedDocument()
        doc.full_text = f"지원하지 않는 파일 형식입니다: {ext}. .hwpx, .pptx, .hwp 파일을 업로드해주세요."
        return doc


def extract_content_lines(full_text: str, max_lines: int = 200) -> list[str]:
    """Clean and filter text lines for storage."""
    lines: list[str] = []
    for line in full_text.split("\n"):
        stripped = line.strip()
        if not stripped:
            continue
        # Skip control characters and very short fragments
        if len(stripped) < 2:
            continue
        # Normalize whitespace
        cleaned = re.sub(r"\s+", " ", stripped)
        lines.append(cleaned)
    return lines[:max_lines]

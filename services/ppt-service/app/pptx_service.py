"""python-pptx wrappers — pure functions, no LLM, no I/O."""

from __future__ import annotations

from io import BytesIO
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

from app.models import PptOutline, SlideSummary


def parse_pptx_bytes(data: bytes, file_name: str) -> PptOutline:
    """Read a .pptx file from bytes and return a structural outline."""
    presentation = Presentation(BytesIO(data))
    slides: list[SlideSummary] = []
    for idx, slide in enumerate(presentation.slides, start=1):
        title = None
        text_chunks: list[str] = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                line = "".join(run.text for run in para.runs).strip()
                if not line:
                    continue
                if title is None and shape == slide.shapes.title:
                    title = line
                else:
                    text_chunks.append(line)
        preview = " / ".join(text_chunks)[:200] if text_chunks else None
        slides.append(SlideSummary(index=idx, title=title, text_preview=preview))
    return PptOutline(file_name=file_name, slide_count=len(slides), slides=slides)


def build_pptx(outline_slides: list[dict], topic: str) -> bytes:
    """Build a minimal but valid .pptx from an outline.

    Each entry in `outline_slides` should have:
      - title: str
      - bullets: list[str]
    """
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    blank_layout = presentation.slide_layouts[6]

    for entry in outline_slides:
        slide = presentation.slides.add_slide(blank_layout)
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(1.2))
        title_tf = title_box.text_frame
        title_tf.text = entry.get("title", "")
        for p in title_tf.paragraphs:
            for run in p.runs:
                run.font.size = Pt(32)
                run.font.bold = True

        body_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.7), Inches(5.0))
        body_tf = body_box.text_frame
        body_tf.word_wrap = True
        for i, bullet in enumerate(entry.get("bullets", [])):
            p = body_tf.paragraphs[0] if i == 0 else body_tf.add_paragraph()
            p.text = f"• {bullet}"
            for run in p.runs:
                run.font.size = Pt(20)

    buffer = BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def save_pptx(data: bytes, output_dir: Path, filename: str) -> Path:
    output_dir.mkdir(parents=True, exist_ok=True)
    path = output_dir / filename
    path.write_bytes(data)
    return path
